"""
슬라이드 렌더러 — 파싱된 SlideData를 python-pptx 슬라이드로 변환
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import os

from . import design as D
from .text_utils import add_rich_text, set_paragraph_format
from .parser import SlideData, Element


# ═══════════════════════════════════════
# 공통 컴포넌트
# ═══════════════════════════════════════

def _set_bg(slide, color):
    """슬라이드 배경색"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_bg(slide, prs, color, left=0, top=0, width=None, height=None):
    """배경 사각형"""
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_accent_line(slide, left, top, width=None):
    """포인트 악센트 라인"""
    w = width or D.ACCENT_LINE_WIDTH
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, w, D.ACCENT_LINE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = D.ACCENT
    shape.line.fill.background()


def _add_textbox(slide, left, top, width, height):
    """텍스트 박스 생성 (word_wrap 기본 설정)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox


def _render_title_bar(slide, text, is_dark=False):
    """상단 제목 + 악센트 라인 (일반 슬라이드 공통)"""
    color = D.TEXT_WHITE if is_dark else D.TEXT_DARK
    txBox = _add_textbox(slide, D.MARGIN_LEFT, D.TITLE_TOP,
                         D.CONTENT_WIDTH, D.TITLE_HEIGHT)
    p = txBox.text_frame.paragraphs[0]
    add_rich_text(p, text, font_size=D.TITLE_FONT_SIZE,
                  color=color, bold=True, font_name=D.FONT_TITLE)
    _add_accent_line(slide, D.MARGIN_LEFT, D.ACCENT_LINE_TOP)


def _render_quote_box(slide, text, top=None):
    """인용문 박스 (하단 고정)"""
    y = top or D.QUOTE_BOTTOM_Y
    width = D.CONTENT_WIDTH
    left = D.MARGIN_LEFT

    # 배경
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, y, width, D.QUOTE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = D.QUOTE_BG
    shape.line.fill.background()

    # 왼쪽 바
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, y, D.QUOTE_BAR_WIDTH, D.QUOTE_HEIGHT
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = D.ACCENT
    bar.line.fill.background()

    # 텍스트
    txBox = _add_textbox(slide, left + Inches(0.3), y + Inches(0.12),
                         width - Inches(0.5), D.QUOTE_HEIGHT - Inches(0.24))
    p = txBox.text_frame.paragraphs[0]
    add_rich_text(p, text, font_size=D.QUOTE_FONT_SIZE,
                  color=D.TEXT_DARK, font_name=D.FONT_BODY)


def _render_code_block(slide, lines, left, top, width, height=None):
    """코드 블록"""
    if height is None:
        height = Inches(max(1.0, len(lines) * 0.32 + 0.4))

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = D.CODE_BG
    shape.line.fill.background()

    txBox = _add_textbox(slide, left + D.CODE_PADDING, top + Inches(0.15),
                         width - D.CODE_PADDING * 2, height - Inches(0.3))
    tf = txBox.text_frame

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        add_rich_text(p, line or " ", font_size=D.CODE_FONT_SIZE,
                      color=D.TEXT_CODE, font_name=D.FONT_CODE)
        p.space_after = Pt(D.CODE_FONT_SIZE * (D.CODE_LINE_SPACING - 1))

    return height


def _render_bullets(slide, lines, left, top, width, font_size=None, color=None):
    """불릿 리스트"""
    fs = font_size or D.BULLET_FONT_SIZE
    clr = color or D.TEXT_DARK
    line_count = len(lines)
    height = Inches(max(0.5, line_count * 0.38))

    txBox = _add_textbox(slide, left, top, width, height)
    tf = txBox.text_frame

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet_text = f"•  {line}"
        add_rich_text(p, bullet_text, font_size=fs, color=clr,
                      font_name=D.FONT_BODY)
        p.space_after = Pt(fs * (D.BODY_LINE_SPACING - 1))

    return height


def _render_table(slide, headers, rows, left, top, width):
    """테이블 렌더링"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    if num_cols == 0:
        return Inches(0.5)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width,
        D.TABLE_ROW_HEIGHT * num_rows
    )
    table = table_shape.table

    # 헤더
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = D.TABLE_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        add_rich_text(p, header, font_size=D.TABLE_HEADER_FONT_SIZE,
                      color=D.TEXT_WHITE, bold=True, font_name=D.FONT_BODY)
        p.alignment = PP_ALIGN.CENTER

    # 데이터
    for i, row in enumerate(rows):
        bg = D.TABLE_ROW_LIGHT if i % 2 == 0 else D.TABLE_ROW_WHITE
        for j in range(num_cols):
            cell = table.cell(i + 1, j)
            val = row[j] if j < len(row) else ""
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell_lines = val.split("\n") if "\n" in val else [val]
            for li, line in enumerate(cell_lines):
                if li == 0:
                    p = cell.text_frame.paragraphs[0]
                else:
                    p = cell.text_frame.add_paragraph()
                add_rich_text(p, line, font_size=D.TABLE_BODY_FONT_SIZE,
                              color=D.TEXT_DARK, font_name=D.FONT_BODY)
                p.alignment = align

    return D.TABLE_ROW_HEIGHT * num_rows


def _render_image(slide, img_path, left, top, max_width, max_height=None):
    """이미지 삽입 — 슬라이드 영역에 맞게 비율 유지하며 중앙 배치"""
    from PIL import Image as PILImage

    # 프로젝트 루트 기준 상대 경로 처리
    if not os.path.isabs(img_path):
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        img_path = os.path.join(base_dir, img_path)

    if not os.path.exists(img_path):
        return Inches(0.5)

    # 이미지 원본 크기 확인
    with PILImage.open(img_path) as img:
        img_w, img_h = img.size

    # 최대 영역
    mw = max_width
    mh = max_height or Inches(4.5)

    # 비율 유지 축소
    ratio_w = mw / Emu(int(img_w * 914400 / 96))  # px → EMU (96dpi 기준)
    ratio_h = mh / Emu(int(img_h * 914400 / 96))
    ratio = min(ratio_w, ratio_h, 1.0)  # 원본보다 커지지 않게

    final_w = int(img_w * 914400 / 96 * ratio)
    final_h = int(img_h * 914400 / 96 * ratio)

    # 수평 중앙 정렬
    center_left = left + (mw - final_w) // 2

    slide.shapes.add_picture(img_path, center_left, top,
                             Emu(final_w), Emu(final_h))
    return Emu(final_h)


def _render_checkbox(slide, lines, left, top, width):
    """체크박스 리스트"""
    height = Inches(max(0.5, len(lines) * 0.38))
    txBox = _add_textbox(slide, left, top, width, height)
    tf = txBox.text_frame

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        add_rich_text(p, line, font_size=D.BULLET_FONT_SIZE,
                      color=D.TEXT_DARK, font_name=D.FONT_BODY)
        p.space_after = Pt(D.BULLET_FONT_SIZE * (D.BODY_LINE_SPACING - 1))

    return height


# ═══════════════════════════════════════
# 슬라이드 타입별 렌더러
# ═══════════════════════════════════════

def _render_title_slide(slide, prs, data: SlideData):
    """[제목] 슬라이드 — 챕터 시작, 섹션 구분"""
    _add_shape_bg(slide, prs, D.BG_DARK)

    # 좌측 장식 바 (세로)
    side_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), 0,
        D.SIDE_BAR_WIDTH, prs.slide_height
    )
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = D.ACCENT
    side_bar.line.fill.background()

    # 요소에서 제목, 부제, 설명, 인용문 추출
    headings = [e for e in data.elements if e.type == "heading"]
    quotes = [e for e in data.elements if e.type == "quote"]

    y = D.TITLE_SLIDE_MAIN_TOP

    # 악센트 라인 (제목 위 — 충분한 간격)
    _add_accent_line(slide, Inches(1.5), y - Inches(0.35))

    # H1 = 메인 제목
    if len(headings) >= 1:
        txBox = _add_textbox(slide, Inches(1.5), y,
                             Inches(10), Inches(1.0))
        p = txBox.text_frame.paragraphs[0]
        add_rich_text(p, headings[0].content,
                      font_size=D.TITLE_SLIDE_MAIN_SIZE,
                      color=D.TEXT_WHITE, bold=True, font_name=D.FONT_TITLE)
        y += Inches(1.0)

    # H2 = 서브 제목
    if len(headings) >= 2:
        txBox = _add_textbox(slide, Inches(1.5), y,
                             Inches(10), Inches(0.7))
        p = txBox.text_frame.paragraphs[0]
        add_rich_text(p, headings[1].content,
                      font_size=D.TITLE_SLIDE_SUB_SIZE,
                      color=D.ACCENT, font_name=D.FONT_TITLE)
        y += Inches(0.7)

    # H3 = 설명
    if len(headings) >= 3:
        txBox = _add_textbox(slide, Inches(1.5), y,
                             Inches(10), Inches(0.6))
        p = txBox.text_frame.paragraphs[0]
        add_rich_text(p, headings[2].content,
                      font_size=D.TITLE_SLIDE_DESC_SIZE,
                      color=D.TEXT_GRAY, font_name=D.FONT_TITLE)
        y += Inches(0.7)

    # 인용문 (하단)
    if quotes:
        _render_quote_box(slide, quotes[-1].content)

    # 나머지 텍스트 요소 (대상, 환경 등)
    other_texts = [e for e in data.elements
                   if e.type in ("text", "bold_text", "bullets")]
    if other_texts:
        y_text = Inches(5.2)
        for elem in other_texts:
            if elem.type == "bullets":
                _render_bullets(slide, elem.lines, Inches(1.5), y_text,
                                Inches(10), font_size=16,
                                color=D.TEXT_LIGHT_GRAY)
                y_text += Inches(len(elem.lines) * 0.35)
            else:
                txBox = _add_textbox(slide, Inches(1.5), y_text,
                                     Inches(10), Inches(0.4))
                p = txBox.text_frame.paragraphs[0]
                add_rich_text(p, elem.content, font_size=16,
                              color=D.TEXT_LIGHT_GRAY, font_name=D.FONT_BODY)
                y_text += Inches(0.4)


def _render_content_slide(slide, prs, data: SlideData, is_dark=False):
    """[내용] 슬라이드 — 범용 콘텐츠"""
    bg = D.BG_DARK if is_dark else D.BG_LIGHT
    _add_shape_bg(slide, prs, bg)

    # 첫 번째 H1을 제목으로 사용
    title_text = ""
    elements = list(data.elements)
    for i, e in enumerate(elements):
        if e.type == "heading" and e.level == 1:
            title_text = e.content
            elements.pop(i)
            break

    _render_title_bar(slide, title_text, is_dark=is_dark)

    # 인용문 분리 (하단 고정)
    quotes = [e for e in elements if e.type == "quote"]
    body_elements = [e for e in elements if e.type != "quote"]

    # Y좌표 자동 누적
    y = D.CONTENT_START_Y
    text_color = D.TEXT_WHITE if is_dark else D.TEXT_DARK

    for elem in body_elements:
        if elem.type == "heading":
            fs = D.SUB_HEADING_FONT_SIZE if elem.level == 2 else 16
            txBox = _add_textbox(slide, D.MARGIN_LEFT, y,
                                 D.CONTENT_WIDTH, Inches(0.45))
            p = txBox.text_frame.paragraphs[0]
            add_rich_text(p, elem.content, font_size=fs,
                          color=text_color, bold=True, font_name=D.FONT_BODY)
            y += Inches(0.45) + D.ELEMENT_GAP

        elif elem.type == "bold_text":
            txBox = _add_textbox(slide, D.MARGIN_LEFT, y,
                                 D.CONTENT_WIDTH, Inches(0.4))
            p = txBox.text_frame.paragraphs[0]
            add_rich_text(p, elem.content, font_size=D.BODY_FONT_SIZE,
                          color=text_color, bold=True, font_name=D.FONT_BODY)
            y += Inches(0.4) + D.ELEMENT_GAP

        elif elem.type == "text":
            txBox = _add_textbox(slide, D.MARGIN_LEFT, y,
                                 D.CONTENT_WIDTH, Inches(0.35))
            p = txBox.text_frame.paragraphs[0]
            add_rich_text(p, elem.content, font_size=D.BODY_FONT_SIZE,
                          color=text_color, font_name=D.FONT_BODY)
            y += Inches(0.35) + D.ELEMENT_GAP

        elif elem.type == "bullets":
            h = _render_bullets(slide, elem.lines, D.MARGIN_LEFT, y,
                                D.CONTENT_WIDTH, color=text_color)
            y += h + D.ELEMENT_GAP

        elif elem.type == "code_block":
            h = _render_code_block(slide, elem.lines, D.MARGIN_LEFT, y,
                                   D.CONTENT_WIDTH)
            y += h + D.ELEMENT_GAP

        elif elem.type == "table":
            h = _render_table(slide, elem.headers, elem.rows,
                              Inches(0.8), y, Inches(11.7))
            y += h + D.ELEMENT_GAP

        elif elem.type == "checkbox":
            h = _render_checkbox(slide, elem.lines, D.MARGIN_LEFT, y,
                                 D.CONTENT_WIDTH)
            y += h + D.ELEMENT_GAP

        elif elem.type == "image":
            h = _render_image(slide, elem.content, D.MARGIN_LEFT, y,
                              D.CONTENT_WIDTH)
            y += h + D.ELEMENT_GAP

    # 인용문 하단 고정
    if quotes:
        _render_quote_box(slide, quotes[-1].content)


def _render_code_slide(slide, prs, data: SlideData):
    """[코드] 슬라이드 — 코드 블록 중심"""
    # 코드 슬라이드도 content와 비슷하지만 코드 블록이 더 강조됨
    _render_content_slide(slide, prs, data, is_dark=False)


def _render_compare_slide(slide, prs, data: SlideData):
    """[비교] 슬라이드 — 좌우 비교 또는 테이블 비교"""
    # 테이블이 있으면 다크 배경, 없으면 라이트
    has_table = any(e.type == "table" for e in data.elements)
    code_blocks = [e for e in data.elements if e.type == "code_block"]

    if has_table and not code_blocks:
        # 테이블 비교 → 다크 배경
        _render_content_slide(slide, prs, data, is_dark=True)
    elif len(code_blocks) >= 2:
        # 좌우 코드 비교
        _render_side_by_side(slide, prs, data, code_blocks)
    else:
        _render_content_slide(slide, prs, data, is_dark=False)


def _render_side_by_side(slide, prs, data: SlideData, code_blocks):
    """좌우 비교 레이아웃 (코드블록 2개)"""
    _add_shape_bg(slide, prs, D.BG_LIGHT)

    # 제목
    title_text = ""
    for e in data.elements:
        if e.type == "heading" and e.level == 1:
            title_text = e.content
            break
    _render_title_bar(slide, title_text)

    # 코드블록 앞의 볼드 텍스트를 라벨로
    labels = [e for e in data.elements if e.type == "bold_text"]
    quotes = [e for e in data.elements if e.type == "quote"]

    half_width = Inches(5.5)
    left_x = D.MARGIN_LEFT
    right_x = Inches(7)
    y_label = D.CONTENT_START_Y
    y_code = D.CONTENT_START_Y + Inches(0.5)

    # 왼쪽 라벨
    if len(labels) >= 1:
        txBox = _add_textbox(slide, left_x, y_label,
                             half_width, Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        add_rich_text(p, labels[0].content, font_size=18,
                      color=D.ACCENT, bold=True, font_name=D.FONT_BODY)

    # 오른쪽 라벨
    if len(labels) >= 2:
        txBox = _add_textbox(slide, right_x, y_label,
                             half_width, Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        add_rich_text(p, labels[1].content, font_size=18,
                      color=D.ACCENT_BLUE, bold=True, font_name=D.FONT_BODY)

    # 왼쪽 코드
    if len(code_blocks) >= 1:
        _render_code_block(slide, code_blocks[0].lines,
                           left_x, y_code, half_width)

    # 오른쪽 코드
    if len(code_blocks) >= 2:
        _render_code_block(slide, code_blocks[1].lines,
                           right_x, y_code, half_width)

    # 인용문
    if quotes:
        _render_quote_box(slide, quotes[-1].content)


def _render_practice_slide(slide, prs, data: SlideData):
    """[실습] 슬라이드 — 실습 단계 + 체크포인트"""
    _render_content_slide(slide, prs, data, is_dark=False)


# ═══════════════════════════════════════
# 메인 렌더링 함수
# ═══════════════════════════════════════

_TYPE_RENDERERS = {
    "title": _render_title_slide,
    "content": _render_content_slide,
    "code": _render_code_slide,
    "compare": _render_compare_slide,
    "practice": _render_practice_slide,
}


def _add_slide_number(slide, num, total, is_dark=False):
    """슬라이드 번호 (우하단)"""
    color = D.TEXT_GRAY if is_dark else RGBColor(0xBB, 0xBB, 0xBB)
    txBox = _add_textbox(slide, Inches(11.5), Inches(7.0),
                         Inches(1.5), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    add_rich_text(p, f"{num} / {total}", font_size=10,
                  color=color, font_name=D.FONT_BODY)


def render_slides(slides_data, output_path):
    """SlideData 리스트를 PPT 파일로 렌더링

    Args:
        slides_data: list[SlideData]
        output_path: 출력 .pptx 경로
    """
    prs = Presentation()
    prs.slide_width = D.SLIDE_WIDTH
    prs.slide_height = D.SLIDE_HEIGHT

    total = len(slides_data)
    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        renderer = _TYPE_RENDERERS.get(slide_data.slide_type,
                                       _render_content_slide)
        renderer(slide, prs, slide_data)

        # 타이틀 슬라이드 제외하고 슬라이드 번호 추가
        is_dark = slide_data.slide_type in ("title",) or (
            slide_data.slide_type == "compare" and
            any(e.type == "table" for e in slide_data.elements) and
            not any(e.type == "code_block" for e in slide_data.elements)
        )
        if slide_data.slide_type != "title":
            _add_slide_number(slide, idx + 1, total, is_dark=is_dark)

    import os
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    return len(slides_data)
